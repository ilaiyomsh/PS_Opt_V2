"""Generate an editable PowerPoint deck for PS_Opt_V2 architecture.

Every box, arrow, and label is a native PowerPoint shape — selectable,
movable, recolorable, retypable directly in PowerPoint / Keynote / Canva.

Run:
    source venv/bin/activate
    python legacy/make_pptx.py
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = Path(__file__).resolve().parent.parent / "architecture_slide.pptx"

# ---------- palette
INK    = RGBColor(0x1F, 0x29, 0x37)
SOFT   = RGBColor(0x6B, 0x72, 0x80)
RULE   = RGBColor(0xD1, 0xD5, 0xDB)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
RED    = RGBColor(0xB9, 0x1C, 0x1C)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
AMBER  = RGBColor(0xB4, 0x53, 0x09)
ACCENT = RGBColor(0xC2, 0x41, 0x0C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
F_BLUE   = RGBColor(0xEF, 0xF6, 0xFF)
F_GREEN  = RGBColor(0xF0, 0xFD, 0xF4)
F_RED    = RGBColor(0xFE, 0xF2, 0xF2)
F_PURPLE = RGBColor(0xF5, 0xF3, 0xFF)
F_AMBER  = RGBColor(0xFF, 0xFB, 0xEB)

# ---------- presentation (16:9)
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_text(shape, text, *, size=12, color=INK, bold=False, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Helvetica Neue"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size if i == 0 else max(8, size - 4))
        r.font.bold = bold and i == 0
        r.font.color.rgb = color if i == 0 else SOFT

def box(x, y, w, h, title, sub=None, *, fill=F_BLUE, edge=BLUE,
        title_sz=14, title_color=INK, edge_w=1.5, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    # rounded corner radius
    if rounded:
        s.adjustments[0] = 0.10
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = edge
    s.line.width = Pt(edge_w)
    s.shadow.inherit = False
    text = title if not sub else f"{title}\n{sub}"
    add_text(s, text, size=title_sz, color=title_color, bold=True)
    return s

def cluster(x, y, w, h, label, edge):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = 0.04
    s.fill.background()
    s.line.color.rgb = edge
    s.line.width = Pt(1.0)
    # dashed
    ln = s.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")
    s.shadow.inherit = False
    # label in a separate textbox at top-left
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05),
                                  Inches(w - 0.3), Inches(0.3))
    add_text(tb, label, size=10, color=edge, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

def label(x, y, w, h, text, *, size=10, color=SOFT, bold=False,
          align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    add_text(tb, text, size=size, color=color, bold=bold, align=align,
             anchor=MSO_ANCHOR.MIDDLE)

def connect(x0, y0, x1, y1, *, color=INK, weight=1.4, dashed=False,
            arrow=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x0), Inches(y0),
                                   Inches(x1), Inches(y1))
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    ln = c.line._get_or_add_ln()
    if dashed:
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return c

# ============================================================ HEADER
label(0.4, 0.20, 12.5, 0.55,
      "PS_Opt_V2", size=26, color=INK, bold=True, align=PP_ALIGN.LEFT)
label(0.4, 0.70, 12.5, 0.35,
      "Silicon PIN phase-shifter optimization — system architecture",
      size=12, color=SOFT, align=PP_ALIGN.LEFT)

# header rule
connect(0.4, 1.10, 12.93, 1.10, color=RULE, weight=0.75, arrow=False)

# ============================================================ CLUSTERS
cluster(0.40, 1.25, 3.55, 5.55,  "CONFIG",  AMBER)
cluster(4.10, 1.25, 4.95, 5.55,  "PER-SIMULATION PIPELINE  ·  run_simulation.run_row", GREEN)
cluster(9.20, 1.25, 3.73, 5.55,  "OPTIMIZATION  ·  STORAGE", PURPLE)

# ============================================================ LEFT COLUMN
# config
box(0.65, 1.65, 3.05, 0.65, "config.py",
    "bounds · weights · targets · flags",
    fill=WHITE, edge=AMBER, title_sz=12)

# main.py
box(0.65, 2.65, 3.05, 0.75, "main.py",
    "stage controller  ·  LHS › sims › BO",
    fill=WHITE, edge=BLUE, title_sz=13)

# LHS
box(0.65, 3.75, 3.05, 0.70, "LHS.py",
    "Latin Hypercube  ·  smt",
    fill=WHITE, edge=GREEN, title_sz=12)

# params.csv
box(0.65, 4.85, 3.05, 0.70, "params.csv",
    "LHS-generated inputs",
    fill=WHITE, edge=PURPLE, title_sz=12)

# arrows in left column
connect(2.18, 2.30, 2.18, 2.65, color=AMBER, weight=1.2, dashed=True)
connect(2.18, 3.40, 2.18, 3.75, color=BLUE,  weight=1.4)
connect(2.18, 4.45, 2.18, 4.85, color=GREEN, weight=1.4)

# ============================================================ CENTER PIPELINE
cx_p_left = 4.35
pipe_w    = 4.45
pipe_h    = 0.65
stages = [
    ("snap_params_dict",                "discrete grid  ·  h_si  ·  10 nm",            GREEN),
    ("sim_handler.run_full_simulation", "CHARGE  ›  FDE   via lumapi",                 GREEN),
    ("data_processor",                  "process_charge_data  ·  C(V)",                GREEN),
    ("data_processor",                  "process_optical_data  ·  Δn_eff  α  Δφ  Vπ",  GREEN),
    ("cost.py  ·  FoM",                 "valid: w·(α/T)² + w·(VπL/T)²    fail: C_BASE+β(π−Δφ)²", GREEN),
]
pipe_ys = [1.65, 2.55, 3.45, 4.35, 5.25]
for (t, s, c), y in zip(stages, pipe_ys):
    box(cx_p_left, y, pipe_w, pipe_h, t, s,
        fill=WHITE, edge=c, title_sz=12)

# down arrows
for y in pipe_ys[:-1]:
    connect(cx_p_left + pipe_w / 2, y + pipe_h,
            cx_p_left + pipe_w / 2, y + 0.90,
            color=GREEN, weight=1.3)

# params.csv → snap (cross-cluster)
connect(3.70, 5.20, 4.35, 1.95,
        color=PURPLE, weight=1.2, dashed=True)
label(3.85, 3.45, 0.6, 0.25, "rows", size=9, color=SOFT, align=PP_ALIGN.LEFT)

# ============================================================ RIGHT COLUMN
# Lumerical (external)
box(9.40, 2.55, 3.30, 0.65, "Lumerical",
    "CHARGE + FDE   (lumapi)",
    fill=F_RED, edge=RED, title_sz=12)
# bidirectional arrows pipeline <-> Lumerical
connect(cx_p_left + pipe_w, 2.80, 9.40, 2.80, color=RED, weight=1.3)
connect(9.40, 2.95, cx_p_left + pipe_w, 2.95, color=RED, weight=1.3)

# BO.py
box(9.40, 1.55, 3.30, 0.80, "BO.py",
    "GP + UCB  ·  norm params  ·  −log(cost)  ·  kappa decay",
    fill=WHITE, edge=ACCENT, title_sz=12)

# result.csv
box(9.40, 3.75, 3.30, 0.65, "result.csv",
    "minimal  ·  BO-facing",
    fill=WHITE, edge=PURPLE, title_sz=12)

# result_full + raw + errors
box(9.40, 4.65, 3.30, 0.65, "result_full.csv  ·  raw/  ·  errors.csv",
    "full record  ·  per-sim sweeps  ·  failure log",
    fill=WHITE, edge=PURPLE, title_sz=10)

# results_archive
box(9.40, 5.65, 3.30, 0.60, "results_archive/",
    "warm-start  ·  prepare_initial_data.py",
    fill=WHITE, edge=SOFT, title_sz=11)

# cost → result.csv
connect(cx_p_left + pipe_w, 5.55, 9.40, 4.00, color=PURPLE, weight=1.3)

# result.csv → BO (train + register)
connect(11.05, 3.75, 11.05, 2.35, color=ACCENT, weight=1.5)
label(11.20, 3.00, 1.6, 0.30, "train + register",
      size=10, color=ACCENT, bold=True, align=PP_ALIGN.LEFT)

# archive → result.csv warm-start (dashed)
connect(11.05, 5.65, 11.05, 4.40, color=SOFT, weight=1.0, dashed=True)
label(11.20, 4.95, 1.6, 0.30, "warm-start",
      size=9, color=SOFT, align=PP_ALIGN.LEFT)

# BO → snap : control-loop (top route)
# from BO top-left, up to header rule area, across, down into snap.py
y_top = 1.30
# segment 1: up from BO
connect(9.55, 1.55, 9.55, y_top, color=ACCENT, weight=1.5, arrow=False)
# segment 2: across
connect(9.55, y_top, 6.57, y_top, color=ACCENT, weight=1.5, arrow=False)
# segment 3: down into snap (with arrowhead)
connect(6.57, y_top, 6.57, 1.65, color=ACCENT, weight=1.5)
label(7.30, y_top - 0.30, 2.5, 0.30, "suggest next x",
      size=10, color=ACCENT, bold=True, align=PP_ALIGN.LEFT)

# ============================================================ LEGEND STRIP
ly = 7.05
legend_items = [
    ("processing", GREEN),
    ("BO loop",    ACCENT),
    ("external",   RED),
    ("data store", PURPLE),
    ("config",     AMBER),
]
x = 0.40
label(x, ly, 0.7, 0.30, "Legend:", size=10, color=SOFT, bold=True)
x = 1.10
for txt, col in legend_items:
    sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(ly + 0.08),
                                Inches(0.22), Inches(0.18))
    sw.fill.solid(); sw.fill.fore_color.rgb = col
    sw.line.fill.background()
    sw.shadow.inherit = False
    label(x + 0.30, ly, 1.5, 0.30, txt, size=10, color=INK,
          align=PP_ALIGN.LEFT)
    x += 1.55

# ============================================================ SAVE
prs.save(OUT)
print(f"Wrote: {OUT}")
